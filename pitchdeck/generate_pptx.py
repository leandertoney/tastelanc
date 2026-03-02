#!/usr/bin/env python3
"""
TasteLanc Pitch Deck Generator
Generates TasteLanc_PitchDeck.pptx from codebase-verified data.

Usage:
    pip install python-pptx Pillow
    python generate_pptx.py

All factual claims sourced from the TasteLanc repository.
See README.md for the complete source map.
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw

# ---------------------------------------------------------------------------
# Brand colors (from apps/mobile/src/constants/colors.ts)
# ---------------------------------------------------------------------------
BG_COLOR = RGBColor(0x1A, 0x1A, 0x1A)       # #1A1A1A deep charcoal
ACCENT = RGBColor(0xA4, 0x1E, 0x22)          # #A41E22 brand red
WHITE = RGBColor(0xFF, 0xFF, 0xFF)            # #FFFFFF
MUTED = RGBColor(0x99, 0x99, 0x99)            # muted text
CARD_BG = RGBColor(0x25, 0x25, 0x25)          # #252525 card bg
GOLD = RGBColor(0xFF, 0xD7, 0x00)             # #FFD700 premium gold
SURFACE = RGBColor(0x1E, 0x1E, 0x1E)          # #1E1E1E surface

# Slide dimensions (16:9 widescreen)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Paths
REPO_ROOT = Path(__file__).resolve().parent.parent
ASSETS_DIR = REPO_ROOT / "apps" / "mobile" / "assets"
WEB_PUBLIC = REPO_ROOT / "apps" / "web" / "public" / "images"
SCREENSHOTS_DIR = ASSETS_DIR / "app-store-screenshots"
DESKTOP_ASSETS = Path.home() / "Desktop" / "TasteLanc Assets"
FRAMES_DIR = Path(__file__).resolve().parent / "frames"
OUTPUT_PATH = Path(__file__).resolve().parent / "TasteLanc_PitchDeck.pptx"


def set_slide_bg(slide, color=BG_COLOR):
    """Set slide background to solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=16, color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(6), space_after=Pt(2),
                  font_name="Calibri"):
    """Add a new paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_accent_bar(slide, left, top, width=Inches(1.5), height=Pt(4)):
    """Add a thin accent color bar (decorative element)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, color=CARD_BG):
    """Add a rounded-corner card background."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Reduce corner rounding
    shape.adjustments[0] = 0.05
    return shape


def try_add_image(slide, image_path, left, top, width=None, height=None):
    """Try to add an image; return True if successful."""
    if image_path.exists():
        kwargs = {"image_file": str(image_path), "left": left, "top": top}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        slide.shapes.add_picture(**kwargs)
        return True
    return False


# ---------------------------------------------------------------------------
# Device frame generator
# ---------------------------------------------------------------------------

def create_device_frame(screenshot_path: Path, output_path: Path):
    """Wrap a screenshot in an iPhone-style device frame using Pillow."""
    ss = Image.open(screenshot_path).convert("RGBA")
    ss_w, ss_h = ss.size  # 1242x2688 for 6.5" screenshots

    # Frame dimensions
    bezel = int(ss_w * 0.04)        # 4% border width (~50px)
    top_bar = int(ss_h * 0.015)     # top chin
    bottom_bar = int(ss_h * 0.015)  # bottom chin
    corner_radius = int(ss_w * 0.08)  # rounded corners

    frame_w = ss_w + bezel * 2
    frame_h = ss_h + bezel * 2 + top_bar + bottom_bar

    # Create frame canvas
    frame = Image.new("RGBA", (frame_w, frame_h), (0, 0, 0, 0))
    draw = ImageDraw.Draw(frame)

    # Draw outer rounded rectangle (device body - dark gray)
    draw.rounded_rectangle(
        [0, 0, frame_w - 1, frame_h - 1],
        radius=corner_radius,
        fill=(30, 30, 30, 255),  # #1E1E1E device body
        outline=(60, 60, 60, 255),  # subtle border
        width=2,
    )

    # Draw inner rounded rectangle (screen area - slightly inset)
    screen_x = bezel
    screen_y = bezel + top_bar
    screen_corner = int(corner_radius * 0.6)
    draw.rounded_rectangle(
        [screen_x - 1, screen_y - 1, screen_x + ss_w, screen_y + ss_h],
        radius=screen_corner,
        fill=(0, 0, 0, 255),
    )

    # Draw Dynamic Island (top notch indicator)
    island_w = int(ss_w * 0.25)
    island_h = int(ss_h * 0.012)
    island_x = (frame_w - island_w) // 2
    island_y = bezel + top_bar + int(ss_h * 0.008)
    island_radius = island_h // 2
    draw.rounded_rectangle(
        [island_x, island_y, island_x + island_w, island_y + island_h],
        radius=island_radius,
        fill=(20, 20, 20, 255),
    )

    # Paste screenshot onto frame
    frame.paste(ss, (screen_x, screen_y))

    # Re-draw the Dynamic Island ON TOP of the screenshot (like real iPhone)
    draw2 = ImageDraw.Draw(frame)
    draw2.rounded_rectangle(
        [island_x, island_y, island_x + island_w, island_y + island_h],
        radius=island_radius,
        fill=(20, 20, 20, 255),
    )

    # Save
    frame.save(str(output_path), "PNG")
    return output_path


def generate_all_device_frames():
    """Generate device-framed versions of all 6 screenshots."""
    FRAMES_DIR.mkdir(exist_ok=True)

    screenshots = [
        "6.5_01_home.png",
        "6.5_02_search.png",
        "6.5_03_rosie_ai.png",
        "6.5_04_happy_hours.png",
        "6.5_05_detail.png",
        "6.5_06_voting.png",
    ]

    framed_paths = []
    for filename in screenshots:
        src = SCREENSHOTS_DIR / filename
        dst = FRAMES_DIR / f"framed_{filename}"
        if src.exists():
            create_device_frame(src, dst)
            framed_paths.append(dst)
            print(f"  \u2713 Frame: {filename}")
        else:
            framed_paths.append(None)
            print(f"  \u2717 Missing: {filename}")

    return framed_paths


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def build_slide_01_cover(prs):
    """Slide 1: Cover"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide)

    # Logo — constrain HEIGHT to 2.5" (logo is 1024x1024 square)
    logo_path = WEB_PUBLIC / "tastelanc_new_dark.png"
    if logo_path.exists():
        logo_h = Inches(2.5)
        logo_w = Inches(2.5)  # square
        logo_x = (SLIDE_WIDTH - logo_w) // 2  # center horizontally
        try_add_image(slide, logo_path, logo_x, Inches(0.6), height=logo_h)
    else:
        add_textbox(slide, Inches(2), Inches(1.0), Inches(9), Inches(1.5),
                    "TasteLanc", font_size=72, color=WHITE, bold=True,
                    alignment=PP_ALIGN.CENTER)

    # Accent bar — well below the logo (logo ends at Y=3.1)
    add_accent_bar(slide, Inches(5.4), Inches(3.5), width=Inches(2.5))

    # Tagline — below accent bar
    add_textbox(slide, Inches(2), Inches(3.9), Inches(9), Inches(0.8),
                "Eat. Drink. Experience.", font_size=40, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, Inches(2), Inches(4.9), Inches(9), Inches(0.6),
                "Lancaster's go-to for what's happening now.",
                font_size=22, color=MUTED, alignment=PP_ALIGN.CENTER)

    # Footer
    add_textbox(slide, Inches(2), Inches(6.5), Inches(9), Inches(0.4),
                "tastelanc.com", font_size=14, color=MUTED,
                alignment=PP_ALIGN.CENTER)


def build_slide_02_oneliner(prs):
    """Slide 2: One-Liner / What Is TasteLanc?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.5),
                "WHAT IS TASTELANC?", font_size=14, color=ACCENT, bold=True)

    add_accent_bar(slide, Inches(0.8), Inches(1.15), width=Inches(2))

    add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(1.2),
                "TasteLanc is the hyper-local discovery platform for restaurants, "
                "happy hours, events, and nightlife \u2014 powered by AI and community voting.",
                font_size=30, color=WHITE, bold=True)

    bullets = [
        "Mobile app (iOS, live on App Store) + responsive web platform",
        "AI concierge (\u201cRosie\u201d) gives personalized recommendations from real local data",
        "Community voting crowns \u201cBest Of\u201d winners across 8 categories each month",
        "Restaurant owners get a full SaaS dashboard (menus, events, specials, analytics)",
        "Multi-market architecture live in 2 Pennsylvania markets",
    ]

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"\u2022  {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.space_before = Pt(8)


def build_slide_03_problem(prs):
    """Slide 3: The Problem"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.5),
                "THE PROBLEM", font_size=14, color=ACCENT, bold=True)

    add_accent_bar(slide, Inches(0.8), Inches(1.15), width=Inches(2))

    # "Sound familiar?" header
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.7),
                "Sound familiar?", font_size=28, color=MUTED, bold=False)

    # Three pain point cards
    pain_points = [
        ("Happy hours scattered across apps.",
         "No single source for what\u2019s on right now."),
        ("Events you hear about too late.",
         "By the time you find out, it\u2019s sold out or over."),
        ("Deals buried in your feed.",
         "Social media algorithms hide the local spots that matter."),
    ]

    for i, (title, desc) in enumerate(pain_points):
        left = Inches(0.8 + i * 4.0)
        top = Inches(2.6)
        card = add_card(slide, left, top, Inches(3.6), Inches(2.2))

        add_textbox(slide, left + Inches(0.3), top + Inches(0.3),
                    Inches(3.0), Inches(0.8),
                    title, font_size=18, color=WHITE, bold=True)

        add_textbox(slide, left + Inches(0.3), top + Inches(1.2),
                    Inches(3.0), Inches(0.8),
                    desc, font_size=14, color=MUTED)

    # Bottom context
    add_textbox(slide, Inches(0.8), Inches(5.3), Inches(11), Inches(1.2),
                "There is no Yelp for \u201cwhat\u2019s happening tonight.\u201d Google Maps "
                "shows hours and reviews \u2014 not happy hour specials, live music lineups, "
                "or weekly deals. Locals piece together Instagram stories, Facebook events, "
                "and word-of-mouth. Visitors are completely lost.",
                font_size=16, color=MUTED)


def build_slide_04_solution(prs):
    """Slide 4: The Solution"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.5),
                "THE SOLUTION", font_size=14, color=ACCENT, bold=True)

    add_accent_bar(slide, Inches(0.8), Inches(1.15), width=Inches(2))

    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.7),
                "The single source of truth for your local food & nightlife scene.",
                font_size=26, color=WHITE, bold=True)

    features = [
        ("AI Concierge (Rosie)", "Personalized recs from real restaurant data"),
        ("Happy Hours", "Real-time deals with item-level pricing"),
        ("Events Calendar", "Live music, trivia, comedy, DJ nights & more"),
        ("Community Voting", "8 monthly categories, 4 votes/user/month"),
        ("Rewards & Loyalty", "Points for check-ins, reviews, photos, referrals"),
        ("Itinerary Builder", "AI-generated day plans with walk-time estimates"),
        ("Geofence Notifications", "Neighborhood-entry alerts & check-in prompts"),
        ("Map Search", "Clustered pins, category filters, distance sorting"),
        ("Weekly Specials", "Owner-managed deals updated from the dashboard"),
    ]

    # 3x3 grid
    for i, (name, desc) in enumerate(features):
        col = i % 3
        row = i // 3
        left = Inches(0.8 + col * 4.0)
        top = Inches(2.5 + row * 1.6)

        add_card(slide, left, top, Inches(3.6), Inches(1.3))
        add_textbox(slide, left + Inches(0.25), top + Inches(0.15),
                    Inches(3.1), Inches(0.5),
                    name, font_size=15, color=ACCENT, bold=True)
        add_textbox(slide, left + Inches(0.25), top + Inches(0.65),
                    Inches(3.1), Inches(0.5),
                    desc, font_size=13, color=MUTED)


def build_slide_05_product(prs, framed_paths=None):
    """Slide 5: Product Demo / Screenshots with iPhone device frames."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.5),
                "PRODUCT", font_size=14, color=ACCENT, bold=True)

    add_accent_bar(slide, Inches(0.8), Inches(1.15), width=Inches(2))

    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
                "Live on the iOS App Store", font_size=24, color=WHITE, bold=True)

    # Screenshot labels
    labels = ["Home Feed", "Map Search", "Rosie AI", "Happy Hours", "Restaurant Detail", "Voting"]
    raw_filenames = [
        "6.5_01_home.png", "6.5_02_search.png", "6.5_03_rosie_ai.png",
        "6.5_04_happy_hours.png", "6.5_05_detail.png", "6.5_06_voting.png",
    ]

    y_top = Inches(2.1)
    img_height = Inches(4.6)
    # 6 phones across a 13.333" slide with even spacing
    phone_width_approx = Inches(2.0)  # framed image is slightly wider than screenshot
    total_width = phone_width_approx * 6
    gap = (SLIDE_WIDTH - total_width - Inches(1.0)) / 5  # distribute remaining space
    x_start = Inches(0.5)

    for i, label in enumerate(labels):
        x = x_start + i * (phone_width_approx + gap)

        # Try framed image first, fall back to raw screenshot
        img_path = None
        if framed_paths and i < len(framed_paths) and framed_paths[i] and framed_paths[i].exists():
            img_path = framed_paths[i]
        else:
            raw_path = SCREENSHOTS_DIR / raw_filenames[i]
            if raw_path.exists():
                img_path = raw_path

        if img_path:
            try:
                slide.shapes.add_picture(str(img_path), x, y_top, height=img_height)
            except Exception:
                add_card(slide, x, y_top, Inches(1.8), img_height)
                add_textbox(slide, x + Inches(0.1), y_top + Inches(2.0),
                            Inches(1.6), Inches(0.5),
                            f"[{label}]", font_size=11, color=MUTED,
                            alignment=PP_ALIGN.CENTER)
        else:
            add_card(slide, x, y_top, Inches(1.8), img_height)
            add_textbox(slide, x + Inches(0.1), y_top + Inches(2.0),
                        Inches(1.6), Inches(0.5),
                        f"Insert: {label}", font_size=11, color=MUTED,
                        alignment=PP_ALIGN.CENTER)

        # Label below
        add_textbox(slide, x, y_top + img_height + Inches(0.1),
                    Inches(2.0), Inches(0.3),
                    label, font_size=12, color=MUTED,
                    alignment=PP_ALIGN.CENTER)


def build_slide_06_how_it_works(prs):
    """Slide 6: How It Works"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.5),
                "HOW IT WORKS", font_size=14, color=ACCENT, bold=True)

    add_accent_bar(slide, Inches(0.8), Inches(1.15), width=Inches(2))

    # Three-sided platform
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
                "Three-Sided Platform", font_size=24, color=WHITE, bold=True)

    sides = [
        ("CONSUMERS", "Free App", [
            "Discover restaurants & nightlife",
            "Vote in 8 monthly categories",
            "Earn rewards via check-ins",
            "Chat with Rosie AI concierge",
        ]),
        ("RESTAURANT OWNERS", "SaaS Dashboard", [
            "Manage menus & photos",
            "Post events & entertainment",
            "Set happy hours & specials",
            "Track analytics & insights",
        ]),
        ("SALES REPS", "CRM + Commission", [
            "Full lead management CRM",
            "AI-generated outreach emails",
            "15\u201320% commission structure",
            "Multi-restaurant cart checkout",
        ]),
    ]

    for i, (title, subtitle, bullets) in enumerate(sides):
        left = Inches(0.6 + i * 4.1)
        top = Inches(2.3)
        add_card(slide, left, top, Inches(3.8), Inches(3.0))

        add_textbox(slide, left + Inches(0.3), top + Inches(0.2),
                    Inches(3.2), Inches(0.4),
                    title, font_size=16, color=ACCENT, bold=True)
        add_textbox(slide, left + Inches(0.3), top + Inches(0.55),
                    Inches(3.2), Inches(0.3),
                    subtitle, font_size=12, color=MUTED)

        txBox = slide.shapes.add_textbox(
            left + Inches(0.3), top + Inches(1.0), Inches(3.2), Inches(1.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        for j, b in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = f"\u2022  {b}"
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"
            p.space_before = Pt(4)

    # Data pipeline
    add_textbox(slide, Inches(0.8), Inches(5.6), Inches(11), Inches(0.4),
                "Data Pipeline", font_size=16, color=WHITE, bold=True)

    pipeline = "Ingest (OutScraper CSV)  \u2192  Enrich (Google Places API)  \u2192  Activate (Sales Onboarding)  \u2192  Engage (Push / AI Blog / Email)"
    add_textbox(slide, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.4),
                pipeline, font_size=14, color=MUTED)


def build_slide_07_business_model(prs):
    """Slide 7: Business Model"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.5),
                "BUSINESS MODEL", font_size=14, color=ACCENT, bold=True)

    add_accent_bar(slide, Inches(0.8), Inches(1.15), width=Inches(2))

    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
                "Five Revenue Streams", font_size=24, color=WHITE, bold=True)

    streams = [
        ("1", "Restaurant SaaS", "$250\u2013$1,100/term", "3 tiers \u00d7 3 billing durations via Stripe", "LIVE"),
        ("2", "Consumer Premium", "$1.99\u2013$29/yr", "2.5x rewards, ad-free, premium voting", "LIVE*"),
        ("3", "Self-Promoter", "$50/month", "DJs, musicians, performers", "LIVE"),
        ("4", "Sponsored Ads", "TBD pricing", "Impression/click tracked, admin-managed", "LIVE"),
        ("5", "Volume Discounts", "10\u201320% off", "Multi-location restaurant groups", "LIVE"),
    ]

    for i, (num, name, price, desc, status) in enumerate(streams):
        top = Inches(2.3 + i * 0.95)
        add_card(slide, Inches(0.8), top, Inches(11.5), Inches(0.8))

        # Number badge
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        Inches(1.1), top + Inches(0.15),
                                        Inches(0.5), Inches(0.5))
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT
        badge.line.fill.background()
        badge.text_frame.paragraphs[0].text = num
        badge.text_frame.paragraphs[0].font.size = Pt(16)
        badge.text_frame.paragraphs[0].font.color.rgb = WHITE
        badge.text_frame.paragraphs[0].font.bold = True
        badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_textbox(slide, Inches(1.8), top + Inches(0.1), Inches(3), Inches(0.35),
                    name, font_size=16, color=WHITE, bold=True)
        add_textbox(slide, Inches(1.8), top + Inches(0.42), Inches(5), Inches(0.3),
                    desc, font_size=12, color=MUTED)
        add_textbox(slide, Inches(8.5), top + Inches(0.1), Inches(2), Inches(0.35),
                    price, font_size=16, color=GOLD, bold=True)
        add_textbox(slide, Inches(10.8), top + Inches(0.15), Inches(1.3), Inches(0.3),
                    status, font_size=11, color=RGBColor(0x34, 0xC7, 0x59), bold=True)

    # Footnote
    add_textbox(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.3),
                "* Consumer premium currently free for all users during launch phase. RevenueCat integrated but unlocked.",
                font_size=11, color=MUTED)


def build_slide_08_pricing(prs):
    """Slide 8: Pricing"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.5),
                "PRICING", font_size=14, color=ACCENT, bold=True)

    add_accent_bar(slide, Inches(0.8), Inches(1.15), width=Inches(2))

    # Restaurant pricing header
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(8), Inches(0.5),
                "Restaurant Plans (B2B SaaS)", font_size=20, color=WHITE, bold=True)

    # Three tier cards
    tiers = [
        ("BASIC", "Free", "Forever", [
            "Hours display",
            "Location on map",
            "Cover photo",
        ], MUTED),
        ("PREMIUM", "$250", "/3mo", [
            "Everything in Basic +",
            "Full menu management",
            "Consumer analytics",
            "Specials & happy hours",
            "Events & entertainment",
            "Push notifications (4/mo)",
        ], WHITE),
        ("ELITE", "$350", "/3mo", [
            "Everything in Premium +",
            "Logo on map",
            "Daily special list",
            "Social media content",
            "Event & music spotlights",
            "Advanced analytics",
        ], GOLD),
    ]

    for i, (name, price, period, features, highlight) in enumerate(tiers):
        left = Inches(0.6 + i * 4.1)
        top = Inches(2.2)
        card_color = CARD_BG if i < 2 else RGBColor(0x2A, 0x22, 0x10)
        add_card(slide, left, top, Inches(3.8), Inches(4.5), color=card_color)

        add_textbox(slide, left + Inches(0.3), top + Inches(0.2),
                    Inches(3.2), Inches(0.35),
                    name, font_size=14, color=highlight, bold=True)

        add_textbox(slide, left + Inches(0.3), top + Inches(0.55),
                    Inches(2), Inches(0.5),
                    price, font_size=36, color=highlight, bold=True)

        add_textbox(slide, left + Inches(2.2), top + Inches(0.7),
                    Inches(1.2), Inches(0.3),
                    period, font_size=14, color=MUTED)

        # Alt prices
        if name == "PREMIUM":
            alt = "$450/6mo  \u00b7  $800/yr"
        elif name == "ELITE":
            alt = "$600/6mo  \u00b7  $1,100/yr"
        else:
            alt = ""
        if alt:
            add_textbox(slide, left + Inches(0.3), top + Inches(1.1),
                        Inches(3.2), Inches(0.3),
                        alt, font_size=12, color=MUTED)

        # Feature list
        feat_top = top + Inches(1.5) if alt else top + Inches(1.2)
        txBox = slide.shapes.add_textbox(
            left + Inches(0.3), feat_top, Inches(3.2), Inches(2.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for j, f in enumerate(features):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = f"\u2713  {f}"
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE if j > 0 or name != "PREMIUM" else ACCENT
            p.font.name = "Calibri"
            p.space_before = Pt(4)

    # Consumer pricing row
    add_textbox(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.3),
                "Consumer: $1.99/mo early access ($4.99 standard)  \u00b7  "
                "Self-Promoter: $50/mo  \u00b7  "
                "Volume: 10\u201320% off for multi-location groups",
                font_size=12, color=MUTED)


def build_slide_09_moat(prs):
    """Slide 9: Defensibility / Moat"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.5),
                "DEFENSIBILITY", font_size=14, color=ACCENT, bold=True)

    add_accent_bar(slide, Inches(0.8), Inches(1.15), width=Inches(2))

    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
                "Five Layers of Moat", font_size=24, color=WHITE, bold=True)

    moats = [
        ("Local Network Effects",
         "Community voting (8 categories, monthly leaderboards) creates engagement "
         "loops. Restaurants promote the app to improve rankings \u2192 more consumers vote \u2192 richer data."),
        ("Proprietary First-Party Data",
         "Custom analytics on Supabase: impressions, clicks, check-ins, area visits, "
         "conversion funnels. No dependency on third-party analytics."),
        ("AI Knowledge Layer",
         "Rosie is trained on market-specific local knowledge \u2014 neighborhood guides, "
         "food scene descriptions, local culture. Can\u2019t be replicated by generic AI."),
        ("Geofence Infrastructure",
         "Radar SDK geofences for restaurants and neighborhoods power automatic "
         "check-in prompts and area-entry notifications."),
        ("Operational Moat",
         "Full sales CRM with AI email generation, commission-based field team, "
         "and multi-step lead nurturing. Expensive for competitors to replicate."),
    ]

    for i, (title, desc) in enumerate(moats):
        top = Inches(2.2 + i * 1.0)
        # Accent dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                      Inches(0.8), top + Inches(0.08),
                                      Inches(0.2), Inches(0.2))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()

        add_textbox(slide, Inches(1.2), top - Inches(0.05), Inches(3), Inches(0.35),
                    title, font_size=16, color=WHITE, bold=True)
        add_textbox(slide, Inches(1.2), top + Inches(0.3), Inches(11), Inches(0.5),
                    desc, font_size=13, color=MUTED)


def build_slide_10_gtm(prs):
    """Slide 10: Go-to-Market"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.5),
                "GO-TO-MARKET", font_size=14, color=ACCENT, bold=True)

    add_accent_bar(slide, Inches(0.8), Inches(1.15), width=Inches(2))

    # Left column: Current (verified)
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
                "Current (Verified in Code)", font_size=20, color=WHITE, bold=True)

    current = [
        "Direct sales team with built-in CRM & AI email generation",
        "15\u201320% commission structure with tiered bonuses",
        "iOS App Store live (ASC ID: 6755852717)",
        "15+ screen onboarding funnel feeds AI personalization",
        "Push notifications via Expo + geofence area alerts",
        "AI-generated blog (Rosie\u2019s Blog) + email via Resend",
        "SEO: sitemaps for restaurants, events, happy hours, blog",
        "Referral rewards: 20 pts per referral (highest action value)",
    ]

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(5.5), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(current):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u2022  {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.space_before = Pt(6)

    # Right column: Strategy (proposal)
    add_textbox(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.5),
                "Strategy (Proposal)", font_size=20, color=MUTED, bold=True)

    proposed = [
        "Influencer partnerships with local food bloggers",
        "Restaurant-as-distribution: co-branded table cards & QR codes",
        "University campus ambassador program",
        "Local event sponsorships for brand visibility",
    ]

    txBox2 = slide.shapes.add_textbox(Inches(7.0), Inches(2.1), Inches(5.5), Inches(3.0))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(proposed):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f"\u2022  {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = MUTED
        p.font.name = "Calibri"
        p.space_before = Pt(6)

    add_textbox(slide, Inches(7.0), Inches(4.0), Inches(5.0), Inches(0.3),
                "\u26a0  Proposals \u2014 not yet implemented in code",
                font_size=11, color=RGBColor(0xFF, 0xD6, 0x0A))


def build_slide_11_expansion(prs):
    """Slide 11: Expansion"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.5),
                "EXPANSION", font_size=14, color=ACCENT, bold=True)

    add_accent_bar(slide, Inches(0.8), Inches(1.15), width=Inches(2))

    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
                "Multi-market architecture is live. Expansion is automated.",
                font_size=22, color=WHITE, bold=True)

    # Live markets
    add_textbox(slide, Inches(0.8), Inches(2.3), Inches(5), Inches(0.4),
                "Live Markets", font_size=16, color=WHITE, bold=True)

    markets = [
        ("Lancaster, PA", "iOS App Store + Web", "LIVE", RGBColor(0x34, 0xC7, 0x59)),
        ("Cumberland County, PA", "iOS TestFlight + Web", "BETA", RGBColor(0xFF, 0xD6, 0x0A)),
    ]

    for i, (name, platform, status, color) in enumerate(markets):
        top = Inches(2.8 + i * 0.5)
        add_textbox(slide, Inches(1.2), top, Inches(3), Inches(0.35),
                    name, font_size=14, color=WHITE)
        add_textbox(slide, Inches(4.5), top, Inches(3), Inches(0.35),
                    platform, font_size=12, color=MUTED)
        add_textbox(slide, Inches(7.5), top, Inches(1.5), Inches(0.35),
                    status, font_size=12, color=color, bold=True)

    # AI Expansion Agent
    add_textbox(slide, Inches(0.8), Inches(4.2), Inches(11), Inches(0.4),
                "AI Expansion Agent (Runs Every 2 Hours)", font_size=16, color=WHITE, bold=True)

    pipeline_steps = [
        ("Suggest", "Proposes cities if pipeline < 20 entries"),
        ("Research", "Census, BEA, OpenStreetMap, College Scorecard data"),
        ("Score", "Weighted model: dining, population, competition, college, income, tourism"),
        ("Brand", "AI generates 3 proposals per city: name, colors, mascot, tagline"),
        ("Jobs", "Creates draft listings: sales rep, market mgr, content, community"),
        ("Review", "Founders vote via email; both must agree to proceed"),
    ]

    # Pipeline flow
    for i, (step, desc) in enumerate(pipeline_steps):
        col = i % 3
        row = i // 3
        left = Inches(0.8 + col * 4.1)
        top = Inches(4.8 + row * 0.9)

        add_textbox(slide, left, top, Inches(0.8), Inches(0.3),
                    step, font_size=13, color=ACCENT, bold=True)
        add_textbox(slide, left + Inches(0.9), top, Inches(3.0), Inches(0.3),
                    desc, font_size=12, color=MUTED)

    # Playbook
    add_textbox(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.3),
                'Playbook: "Add to MARKET_CONFIG \u2192 Seed DB row \u2192 Deploy Netlify site"',
                font_size=12, color=MUTED)


def build_slide_12_traction(prs):
    """Slide 12: Traction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.5),
                "TRACTION", font_size=14, color=ACCENT, bold=True)

    add_accent_bar(slide, Inches(0.8), Inches(1.15), width=Inches(2))

    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
                "What\u2019s Live Today", font_size=24, color=WHITE, bold=True)

    metrics = [
        ("Markets", "2 (Lancaster PA live, Cumberland County PA beta)"),
        ("iOS App", "Published on App Store"),
        ("Web Platform", "tastelanc.com + cumberland.tastelanc.com"),
        ("Restaurant Dashboard", "Full owner portal: analytics, menus, events, happy hours, specials, team"),
        ("Revenue Infrastructure", "Stripe billing with 3 restaurant tiers, consumer subs, self-promoter plans"),
        ("Ad Platform", "featured_ads with impression/click tracking + CTR analytics"),
        ("Sales CRM", "Lead management, AI email generation, commission tracking"),
        ("AI Systems", "Rosie chat, expansion agent (2h cycle), blog gen, email gen, itinerary gen"),
        ("Analytics", "Custom Supabase: page views, impressions, clicks, check-ins, area visits"),
    ]

    for i, (label, value) in enumerate(metrics):
        top = Inches(2.2 + i * 0.52)
        add_textbox(slide, Inches(0.8), top, Inches(3.2), Inches(0.35),
                    label, font_size=14, color=ACCENT, bold=True)
        add_textbox(slide, Inches(4.2), top, Inches(8), Inches(0.35),
                    value, font_size=14, color=WHITE)

    # TBD note
    add_card(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5),
             color=RGBColor(0x2A, 0x22, 0x10))
    add_textbox(slide, Inches(1.1), Inches(6.85), Inches(11), Inches(0.35),
                "User counts, MRR, restaurant signups, and retention tracked in live dashboards. Available upon request.",
                font_size=12, color=GOLD)


def build_slide_13_roadmap(prs):
    """Slide 13: Roadmap"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.5),
                "ROADMAP", font_size=14, color=ACCENT, bold=True)

    add_accent_bar(slide, Inches(0.8), Inches(1.15), width=Inches(2))

    items = [
        ("NEXT", "Home Screen Personalization",
         "RecommendedSection component built, not yet rendered on HomeScreen",
         "IMPLEMENTATION_PLAN.md Phase 3"),
        ("NEXT", "Impression-Feedback Loop",
         "Auto-correct visibility imbalances across restaurant tiers",
         "IMPLEMENTATION_PLAN.md Phase 4"),
        ("PLANNED", "View All Sort/Filter Chips",
         "Enhanced browse experience with sort and filter options",
         "IMPLEMENTATION_PLAN.md Phase 5"),
        ("PLANNED", "Daily Trivia",
         "Reward action type exists (1 pt), marked DEFERRED in code",
         "rewards/points.ts"),
        ("PLANNED", "Android Play Store",
         "App buildable, listing not yet created",
         "brand.ts:25"),
        ("PLANNED", "Self-Serve Restaurant Claiming",
         "Currently sales-driven; self-service not yet built",
         "Not in codebase"),
        ("ONGOING", "Autonomous City Expansion",
         "AI agent actively researching new markets",
         "expansion-agent.ts"),
    ]

    for i, (status, title, desc, source) in enumerate(items):
        top = Inches(1.6 + i * 0.8)

        # Status badge
        badge_color = ACCENT if status == "NEXT" else (
            GOLD if status == "ONGOING" else MUTED)
        add_textbox(slide, Inches(0.8), top, Inches(1.2), Inches(0.3),
                    status, font_size=11, color=badge_color, bold=True)

        add_textbox(slide, Inches(2.2), top, Inches(4), Inches(0.3),
                    title, font_size=15, color=WHITE, bold=True)
        add_textbox(slide, Inches(2.2), top + Inches(0.3), Inches(7), Inches(0.3),
                    desc, font_size=12, color=MUTED)
        add_textbox(slide, Inches(9.5), top + Inches(0.05), Inches(3.5), Inches(0.3),
                    source, font_size=10, color=RGBColor(0x66, 0x66, 0x66))


def build_slide_14_ask(prs):
    """Slide 14: The Ask"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Logo — constrain HEIGHT (logo is 1024x1024 square)
    logo_path = WEB_PUBLIC / "tastelanc_new_dark.png"
    logo_h = Inches(2.0)
    if logo_path.exists():
        logo_x = (SLIDE_WIDTH - logo_h) // 2  # center (square)
        try_add_image(slide, logo_path, logo_x, Inches(0.5), height=logo_h)

    # Accent bar — below logo (logo ends at Y=2.5)
    add_accent_bar(slide, Inches(5.4), Inches(2.8), width=Inches(2.5))

    # CTA text — below accent bar
    add_textbox(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.8),
                "Partner with us to scale the\nhyper-local playbook to new markets.",
                font_size=28, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    # What we're looking for
    asks = [
        "Strategic partners to accelerate market expansion",
        "Investment to scale the sales team and enter new markets",
        "Restaurant group partnerships for anchor tenant acquisition",
    ]

    txBox = slide.shapes.add_textbox(Inches(3), Inches(4.4), Inches(7), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, ask in enumerate(asks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u2022  {ask}"
        p.font.size = Pt(18)
        p.font.color.rgb = MUTED
        p.font.name = "Calibri"
        p.space_before = Pt(8)
        p.alignment = PP_ALIGN.LEFT

    # Contact
    add_textbox(slide, Inches(2), Inches(6.2), Inches(9), Inches(0.3),
                "tastelanc.com  \u00b7  support@tastelanc.com",
                font_size=16, color=MUTED, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(2), Inches(6.7), Inches(9), Inches(0.3),
                "Generated from codebase \u2014 see /pitchdeck/README.md for source map",
                font_size=10, color=RGBColor(0x55, 0x55, 0x55),
                alignment=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Generate device-framed screenshots first
    print("Generating iPhone device frames...")
    framed_paths = generate_all_device_frames()

    print("\nBuilding TasteLanc Pitch Deck...")

    # Build slides in order
    print("  \u2713 Slide: Cover")
    build_slide_01_cover(prs)
    print("  \u2713 Slide: What Is TasteLanc?")
    build_slide_02_oneliner(prs)
    print("  \u2713 Slide: The Problem")
    build_slide_03_problem(prs)
    print("  \u2713 Slide: The Solution")
    build_slide_04_solution(prs)
    print("  \u2713 Slide: Product Screenshots")
    build_slide_05_product(prs, framed_paths=framed_paths)
    print("  \u2713 Slide: How It Works")
    build_slide_06_how_it_works(prs)
    print("  \u2713 Slide: Business Model")
    build_slide_07_business_model(prs)
    print("  \u2713 Slide: Pricing")
    build_slide_08_pricing(prs)
    print("  \u2713 Slide: Defensibility")
    build_slide_09_moat(prs)
    print("  \u2713 Slide: Go-to-Market")
    build_slide_10_gtm(prs)
    print("  \u2713 Slide: Expansion")
    build_slide_11_expansion(prs)
    print("  \u2713 Slide: Traction")
    build_slide_12_traction(prs)
    print("  \u2713 Slide: Roadmap")
    build_slide_13_roadmap(prs)
    print("  \u2713 Slide: The Ask")
    build_slide_14_ask(prs)

    prs.save(str(OUTPUT_PATH))
    print(f"\nSaved to: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
